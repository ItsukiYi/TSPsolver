"""
Fill poster_template.pptx with CS240 project content.
Decorative images/icons are preserved; only text + content images are replaced.
"""
import sys, io, os, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from lxml import etree

TEMPLATE = 'poster_template.pptx'
OUTPUT = 'poster_filled.pptx'

# All image paths relative to project root
PROJECT_ROOT = os.path.dirname(os.path.abspath(__file__))

def full_path(rel):
    return os.path.join(PROJECT_ROOT, rel)

# ── Content images to replace ──
# Map: shape index (0-based) → image file path
IMAGE_REPLACEMENTS = {
    9:  'outputs/visualizations/00_what_is_tsp.png',      # Left: TSP problem viz
    10: 'outputs/city_delivery_500.png',                   # Left lower: city scenario
    13: 'outputs/algorithm_processes.png',                 # Left: algorithm flow
    20: 'outputs/arch_DIFUSCO.png',                        # Middle: DIFUSCO arch
    21: 'outputs/pareto_frontier.png',                     # Middle: results
    22: 'outputs/improvement_strategies.png',              # Middle: results 2
    23: 'outputs/difusco_diffusion_steps.png',             # Right: ablation
    24: 'outputs/arch_DualOpt.png',                        # Middle: DualOpt arch
    25: 'outputs/arch_pipeline.png',                       # Right: pipeline
    26: 'outputs/improvement4_per_instance.png',           # Right: per-instance
}

# ── Decorative shapes to SKIP (indices) ──
DECORATIVE_PICTURES = {4, 5, 18, 29, 31, 33, 38, 41, 43}
DECORATIVE_OTHER = {0, 11, 17, 28, 30, 32, 40}

# ── Text replacements ──
# Map: shape index → new text
TEXT_REPLACEMENTS = {
    1: "Classical Meets Modern: Neural Diffusion for Delivery Route Optimization",

    2: "Pan Qiao\nCS240: Algorithm Design and Analysis · ShanghaiTech University · June 2026",

    6: "Problem & Motivation",  # was "Background & Motivasion"

    3: (
        "The Traveling Salesman Problem (TSP) — given n locations in a metric space, "
        "find the shortest Hamiltonian cycle visiting each exactly once — is a "
        "cornerstone of combinatorial optimization. Despite its simple formulation, "
        "TSP is NP-hard: no polynomial exact algorithm exists unless P=NP.\n\n"
        "TSP drives modern logistics: Meituan, Uber Eats, and Amazon all solve massive "
        "TSP instances daily for last-mile delivery routing. Beyond logistics, "
        "applications span PCB drilling, VLSI chip design, and genome sequencing.\n\n"
        "TSP is algorithmically unique because it (i) is NP-hard yet admits a 1.5× "
        "constant-factor approximation (Christofides); (ii) is simple enough for deep "
        "learning to learn useful heuristics; (iii) is complex enough that pure ML fails "
        "without classical components; and (iv) serves as the canonical benchmark where "
        "new algorithmic ideas are validated before transferring to richer routing problems."
    ),

    12: "Classical Algorithms",  # was "VA-VAE: Vision-Aligned VAE"

    14: "Four classical TSP solvers implemented from scratch in Python/NumPy, spanning the accuracy-speed spectrum.",

    16: (
        "✓ Nearest Neighbor — O(n²) greedy heuristic, 20-30% optimality gap\n"
        "✓ Christofides (1976) — 1.5× approximation guarantee, O(n³) via Blossom MWPM\n"
        "✓ 2-opt Local Search — NumPy-vectorized edge-swap refinement, ~10× speedup\n"
        "✓ LKH3 (Helsgaun 2017) — Gold-standard heuristic, ~0.4% mean optimality gap"
    ),

    7: "Modern Neural Methods",  # was "LightningDiT"

    15: (
        "DIFUSCO (NeurIPS 2023): Casts TSP as categorical denoising diffusion over "
        "graph adjacency matrices. A 12-layer Anisotropic Gated GNN (~5.3M params) "
        "iteratively denoises random Bernoulli noise into a valid tour over 50 steps. "
        "We reproduce this on a single RTX 2060 (original: 8× GPUs) with 3 compatibility "
        "patches for Lightning v2 / CUDA 12.6 / Cython fallback. A pure-PyTorch sparse "
        "GNN replacement was implemented to eliminate torch_sparse dependency.\n\n"
        "DualOpt (AAAI 2025): Dual divide-and-optimize strategy — grid-based plane "
        "partitioning + LKH3 sub-solves, then cascaded neural revisers (REINFORCE-trained, "
        "~710K total params) refine sub-tours at 3 granularities (k=50, 20, 10). "
        "Combines classical decomposition with learned local search."
    ),

    19: "Experimental Results",  # was "Results" — keep similar

    34: "Hybrid classical-neural pipelines outperform either paradigm alone — DIFUSCO → DualOpt achieves best overall results.",

    8: "Key Insights",  # was "Discussion"

    27: (
        "1. LKH3 dominates the speed-quality frontier: 0.4% mean gap, 2.5s for "
        "TSP-1002 — 31× faster than Christofides at n=1000.\n\n"
        "2. DIFUSCO generalizes surprisingly well: trained only on TSP-50, achieves "
        "7.1% gap on TSP-1002 (20× training size). The GNN learns size-agnostic "
        "edge quality representations.\n\n"
        "3. DualOpt excels within training window: near-perfect on berlin52 (0.03% gap) "
        "but degrades beyond n≈100 due to fixed reviser window sizes.\n\n"
        "4. Ablation reveals 2-opt is the dominant quality driver in DIFUSCO: raw "
        "diffusion has ~13% gap; adding 2-opt closes it to ~0.3%. Only 10 inference "
        "steps needed — 50 steps gives minimal extra gain.\n\n"
        "5. Our DIFUSCO → DualOpt hybrid pipeline achieves the best results: −1.67% "
        "vs ground truth on TSP-50, combining diffusion's global pattern capture with "
        "the reviser's learned local optimization."
    ),

    35: (
        "🏆 Best Pipeline (ours): DIFUSCO → DualOpt\n"
        "TSP-50: −1.67% vs GT  |  kroA100: 1.81% gap\n"
        "DualOpt reviser extracts 24% more improvement than 2-opt alone"
    ),

    36: (
        "[1] Sun & Yang. DIFUSCO: Graph-based Diffusion Solvers for Combinatorial "
        "Optimization. NeurIPS, 2023.\n"
        "[2] Zhou et al. DualOpt: A Dual Divide-and-Optimize Algorithm for the "
        "Large-Scale Traveling Salesman Problem. AAAI, 2025.\n"
        "[3] Christofides. Worst-case analysis of a new heuristic for the travelling "
        "salesman problem. Technical Report, 1976.\n"
        "[4] Helsgaun. An Extension of the Lin-Kernighan-Helsgaun TSP Solver. 2017.\n"
        "[5] Kool, van Hoof & Welling. Attention, Learn to Solve Routing Problems! ICLR, 2019."
    ),

    37: "References",  # keep

    39: "Contact",  # was "Want to Talk More?"

    42: "panqiao2025@shanghaitech.edu.cn",
}


def replace_text_preserve_first_run_font(shape, new_text):
    """
    Replace all text in a shape while preserving font properties from the first run.
    Falls back to safe text assignment if no runs exist.
    """
    tf = shape.text_frame

    # Try to capture font from first available run
    sample_font = None
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text.strip():
                sample_font = {
                    'size': run.font.size,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'name': run.font.name,
                }
                # Try to get color
                try:
                    if run.font.color and run.font.color.rgb:
                        sample_font['color'] = run.font.color.rgb
                except:
                    pass
                break
        if sample_font:
            break

    # Clear existing paragraphs
    tf.clear()

    # Split new text into paragraphs
    lines = new_text.split('\n')

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = line

        # Apply captured font
        if sample_font:
            if sample_font.get('size'):
                run.font.size = sample_font['size']
            if sample_font.get('bold') is not None:
                run.font.bold = sample_font['bold']
            if sample_font.get('italic') is not None:
                run.font.italic = sample_font['italic']
            if sample_font.get('name'):
                run.font.name = sample_font['name']
            if sample_font.get('color'):
                try:
                    run.font.color.rgb = sample_font['color']
                except:
                    pass


def replace_image_in_shape(slide, shape, new_image_path):
    """Delete old picture shape and add new one at same position/size."""
    if not os.path.exists(new_image_path):
        print(f"  WARNING: Image not found: {new_image_path}")
        return False

    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    # Remove old shape from XML
    sp = shape._element
    sp.getparent().remove(sp)

    # Add new picture
    try:
        slide.shapes.add_picture(new_image_path, left, top, width, height)
        return True
    except Exception as e:
        print(f"  ERROR adding picture {new_image_path}: {e}")
        return False


def main():
    print(f"Opening template: {TEMPLATE}")
    prs = Presentation(TEMPLATE)
    slide = prs.slides[0]

    # Collect all shapes first (so indices don't shift during modification)
    shapes = list(slide.shapes)
    print(f"Found {len(shapes)} shapes")

    # ── Phase 1: Replace text ──
    print("\n── Phase 1: Text replacement ──")
    text_modified = 0
    for idx, shape in enumerate(shapes):
        if idx in TEXT_REPLACEMENTS:
            new_text = TEXT_REPLACEMENTS[idx]
            old_text = shape.text[:80].replace('\n', ' | ') if hasattr(shape, 'text') and shape.text else '(empty)'
            try:
                replace_text_preserve_first_run_font(shape, new_text)
                print(f"  [{idx}] \"{shape.name}\": OK")
                print(f"       Old: {old_text}")
                print(f"       New: {new_text[:100]}...")
                text_modified += 1
            except Exception as e:
                print(f"  [{idx}] \"{shape.name}\": ERROR - {e}")

    print(f"  Modified {text_modified} text shapes")

    # ── Phase 2: Replace content images ──
    print("\n── Phase 2: Image replacement ──")
    # Re-collect shapes since text modifications may change things
    # Actually, text modifications don't change the shape list.
    # But to be safe, let me work with the XML elements directly.

    shapes_after_text = list(slide.shapes)
    image_replaced = 0
    image_skipped_decorative = 0
    image_missing = 0

    # Process in reverse order so deletions don't affect indices of remaining shapes
    for idx in sorted(IMAGE_REPLACEMENTS.keys(), reverse=True):
        if idx >= len(shapes_after_text):
            print(f"  [{idx}] Index out of range, skipping")
            continue

        shape = shapes_after_text[idx]
        if shape.shape_type != 13:  # PICTURE
            print(f"  [{idx}] \"{shape.name}\" is not a PICTURE (type={shape.shape_type}), skipping")
            continue

        img_path = full_path(IMAGE_REPLACEMENTS[idx])
        print(f"  [{idx}] \"{shape.name}\" → {IMAGE_REPLACEMENTS[idx]}")

        if replace_image_in_shape(slide, shape, img_path):
            image_replaced += 1
        else:
            image_missing += 1

    # Report on decorative images left untouched
    for idx in sorted(DECORATIVE_PICTURES):
        if idx < len(shapes_after_text):
            shape = shapes_after_text[idx]
            if shape.shape_type == 13:  # PICTURE
                image_skipped_decorative += 1

    print(f"\n  Replaced: {image_replaced} content images")
    print(f"  Skipped (decorative): {image_skipped_decorative} images")
    print(f"  Missing/errors: {image_missing}")

    # ── Save ──
    print(f"\n── Saving to {OUTPUT} ──")
    prs.save(OUTPUT)
    print(f"Done! Output saved to {OUTPUT}")

    # Print summary
    print(f"\n=== Summary ===")
    print(f"Text shapes filled: {text_modified}")
    print(f"Content images replaced: {image_replaced}")
    print(f"Decorative images preserved: {image_skipped_decorative}")
    print(f"Output file: {full_path(OUTPUT)}")


if __name__ == '__main__':
    main()
